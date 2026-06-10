"""
Text extractor for PDF, DOCX, PPTX, TXT files.
Returns list of (page_number, text) tuples.
"""
import os
from typing import List, Tuple
import logging

logger = logging.getLogger(__name__)


def extract_text(file_path: str, file_ext: str) -> List[Tuple[int, str]]:
    """
    Extract text from a document file.
    Returns list of (page_number, page_text) tuples.
    """
    ext = file_ext.lower().lstrip(".")

    if ext == "pdf":
        return _extract_pdf(file_path)
    elif ext == "docx":
        return _extract_docx(file_path)
    elif ext == "pptx":
        return _extract_pptx(file_path)
    elif ext == "txt":
        return _extract_txt(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")


def _extract_pdf(file_path: str) -> List[Tuple[int, str]]:
    import fitz  # PyMuPDF
    pages = []
    with fitz.open(file_path) as doc:
        for page_num, page in enumerate(doc, start=1):
            text = page.get_text("text")
            if text.strip():
                pages.append((page_num, text))
    logger.info(f"PDF extracted: {len(pages)} pages from {file_path}")
    return pages


def _extract_docx(file_path: str) -> List[Tuple[int, str]]:
    from docx import Document
    doc = Document(file_path)
    full_text = "\n".join(para.text for para in doc.paragraphs if para.text.strip())
    # DOCX has no natural page breaks — treat as single "page 1"
    return [(1, full_text)] if full_text else []



def _extract_pptx(file_path: str) -> List[Tuple[int, str]]:
    from pptx import Presentation
    prs = Presentation(file_path)
    slides = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            slides.append((slide_num, "\n".join(texts)))
    logger.info(f"PPTX extracted: {len(slides)} slides from {file_path}")
    return slides


def _extract_txt(file_path: str) -> List[Tuple[int, str]]:
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        text = f.read()
    return [(1, text)] if text.strip() else []
